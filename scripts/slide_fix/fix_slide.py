from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path("artifacts/pure-vs-hpe-metro-ransomware.pptx")


def add_textbox(slide, x, y, w, h, text, size, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_card(slide, x, y, w, h, accent_rgb, title, bullets, badge):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(248, 248, 248)
    card.line.color.rgb = RGBColor(215, 215, 215)
    card.line.width = Pt(1)

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_rgb
    accent.line.fill.background()

    badge_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(y + 0.10), Inches(0.42), Inches(0.24))
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = accent_rgb
    badge_shape.line.fill.background()
    add_textbox(slide, x + 0.12, y + 0.10, 0.42, 0.24, badge, 11.5, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)

    add_textbox(slide, x + 0.58, y + 0.09, w - 0.70, 0.26, title, 13.5, True)

    text = "\n".join([f"• {b}" for b in bullets])
    body = add_textbox(slide, x + 0.12, y + 0.38, w - 0.24, h - 0.46, text, 11.5)
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(11.5)
        p.line_spacing = 1.15


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hpe_orange = RGBColor(227, 121, 47)
    pure_green = RGBColor(47, 157, 87)

    margin_x = 0.35
    top_margin = 0.30
    bottom_margin = 0.30
    gutter = 0.35
    content_w = 13.333 - margin_x * 2
    col_w = (content_w - gutter) / 2
    left_x = margin_x
    right_x = margin_x + col_w + gutter

    add_textbox(
        slide,
        margin_x,
        top_margin,
        content_w,
        0.48,
        "Pure vs HPE: Metro Active/Active + Ransomware Resilience",
        32,
        True,
    )
    add_textbox(
        slide,
        margin_x,
        top_margin + 0.47,
        content_w,
        0.28,
        "Visually comparing the solution stack and where complexity shows up for architects and day-2 ops",
        15,
    )

    for x, color, title in [
        (left_x, hpe_orange, "HPE Peer Persistence (portfolio-dependent)"),
        (right_x, pure_green, "Pure ActiveCluster (single platform)"),
    ]:
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.10), Inches(col_w), Inches(0.36))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        add_textbox(slide, x + 0.02, 1.10, col_w - 0.04, 0.36, title, 17, True, RGBColor(255, 255, 255))

    add_textbox(slide, left_x, 1.50, col_w, 0.24, "More decision points: platform → features → licensing → runbooks", 11.5)
    add_textbox(slide, right_x, 1.50, col_w, 0.24, "Same OS + same features across all arrays → simpler design & ops", 11.5)

    headers = ["PLATFORM", "METRO HA (SYNC)", "RANSOMWARE RECOVERY", "DR / VAULT OPTIONS", "OPS & LICENSING"]
    left_cards = [
        ("3P", "3PAR / Primera", ["OS/feature set varies by model/release", "Different lineage and constraints", "Different architecture and tooling"]),
        ("PP", "Remote Copy Groups + Peer Persistence", ["Sync replication via Remote Copy groups", "ATF requires Quorum Witness + policies", "Behavior/ops vary with platform choices"]),
        ("VL", "Snapshots + Virtual Lock (immutability)", ["Locked snapshots cannot be deleted/modified", "Protection may be platform/license-dependent", "Often paired with broader portfolio tools"]),
        ("BK", "Backup / Vault options", ["StoreOnce", "Zerto vault", "3rd-party stacks"]),
        ("OPS", "Ops & licensing", ["More SKUs/entitlements", "Platform-specific runbooks", "More validation effort"]),
    ]
    right_cards = [
        ("OS", "FlashArray (any model) + Purity OS", ["One OS and one feature set across platform", "Hardware choice does not change capabilities", "Consistent behavior across arrays"]),
        ("AC", "ActiveCluster + Pure1 Cloud Mediator", ["Synchronous active/active replication", "Built-in; mediator acts as quorum witness", "Consistent operational behavior"]),
        ("SM", "Snapshots + SafeMode (retention lock)", ["Point-in-time snapshots are read-only", "SafeMode prevents deletion within retention", "Fits directly alongside ActiveCluster"]),
        ("DR", "ActiveDR (async DR)", ["Long-distance DR option", "Same platform and workflows", "Aligned operations model"]),
        ("OPS", "Ops & licensing", ["Fewer SKUs", "One runbook across arrays", "Lower design friction"]),
    ]

    card_gap = 0.18
    first_card_y = 2.10
    card_h = (7.5 - bottom_margin - first_card_y - card_gap * 4 - 0.62) / 5

    for i, header in enumerate(headers):
        y = first_card_y + i * (card_h + card_gap)
        add_textbox(slide, left_x, y - 0.17, col_w, 0.15, header, 13.5, True)
        add_textbox(slide, right_x, y - 0.17, col_w, 0.15, header, 13.5, True)

        l_badge, l_title, l_bullets = left_cards[i]
        r_badge, r_title, r_bullets = right_cards[i]
        add_card(slide, left_x, y, col_w, card_h, hpe_orange, l_title, l_bullets, l_badge)
        add_card(slide, right_x, y, col_w, card_h, pure_green, r_title, r_bullets, r_badge)

    net_y = 7.5 - bottom_margin - 0.52
    net = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(margin_x), Inches(net_y), Inches(content_w), Inches(0.52))
    net.fill.solid()
    net.fill.fore_color.rgb = RGBColor(238, 238, 238)
    net.line.fill.background()
    add_textbox(slide, margin_x + 0.03, net_y + 0.10, content_w - 0.06, 0.30,
                "Net effect: HPE = more platform-driven variability + more components to assemble | "
                "Pure = single platform stack (ActiveCluster + SafeMode) with consistent behavior",
                11.5)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))


if __name__ == "__main__":
    main()
